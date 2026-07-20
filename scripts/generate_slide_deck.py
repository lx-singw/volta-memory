#!/usr/bin/env python3
"""Generate the judge deck from the versioned evaluation artifact.

This script intentionally refuses to parse BENCHMARKS.md or use placeholder
metrics. The Markdown report, slide deck, and submission copy all start from
``backend/eval/artifacts/evaluation-summary.json``.
"""

from __future__ import annotations

import json
import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


REPO_ROOT = Path(__file__).resolve().parents[1]
SUMMARY_PATH = REPO_ROOT / "backend" / "eval" / "artifacts" / "evaluation-summary.json"
OUTPUT_PATH = REPO_ROOT / "docs" / "volta_memory_deck.pptx"

NAVY = RGBColor(4, 14, 30)
SURFACE = RGBColor(13, 33, 56)
SURFACE_ALT = RGBColor(18, 44, 73)
WHITE = RGBColor(248, 250, 252)
MUTED = RGBColor(163, 184, 209)
GOLD = RGBColor(251, 191, 36)
CYAN = RGBColor(34, 211, 238)
GREEN = RGBColor(74, 222, 128)
RED = RGBColor(251, 113, 133)


def load_summary(path: Path) -> dict:
    if not path.exists():
        raise FileNotFoundError(
            f"Missing {path}. Run `python scripts/generate_evaluation_artifacts.py --write-benchmarks` first."
        )
    return json.loads(path.read_text(encoding="utf-8"))


def add_background(slide) -> None:
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = NAVY
    background.line.fill.background()


def add_text(slide, text: str, left: float, top: float, width: float, height: float, *, size: int, color=WHITE, bold=False, align=None) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.size = Pt(size)
    paragraph.font.color.rgb = color
    paragraph.font.bold = bold
    if align is not None:
        paragraph.alignment = align


def add_header(slide, eyebrow: str, title: str) -> None:
    add_text(slide, eyebrow.upper(), 0.8, 0.42, 11.7, 0.28, size=10, color=GOLD, bold=True)
    add_text(slide, title, 0.8, 0.72, 11.7, 0.7, size=30, color=WHITE, bold=True)
    add_text(
        slide,
        "Track 1: MemoryAgent  |  Qwen Cloud Global AI Hackathon",
        0.8,
        7.08,
        11.7,
        0.2,
        size=9,
        color=MUTED,
    )


def add_card(slide, left: float, top: float, width: float, height: float) -> None:
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = SURFACE
    card.line.color.rgb = SURFACE_ALT


def metric_text(system: dict, metric: str) -> str:
    result = system["metrics"][metric]
    if not result["total"]:
        return "N/A"
    return f"{result['ratio']:.1%} ({result['hits']}/{result['total']})"


def add_metric_table(slide, summary: dict) -> None:
    systems = list(summary["systems"].values())
    rows, cols = len(systems) + 1, 6
    shape = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(1.75), Inches(11.75), Inches(3.75))
    table = shape.table
    widths = [2.05, 1.45, 1.55, 1.55, 1.45, 1.7]
    for index, width in enumerate(widths):
        table.columns[index].width = Inches(width)
    headers = ["SYSTEM", "RECALL", "CORRECTION", "FORGETTING", "QUALITY", "RUNS"]
    for column, title in enumerate(headers):
        cell = table.cell(0, column)
        cell.text = title
        cell.fill.solid()
        cell.fill.fore_color.rgb = SURFACE_ALT
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.size = Pt(11)
        paragraph.font.bold = True
        paragraph.font.color.rgb = GOLD
        paragraph.alignment = PP_ALIGN.CENTER
    for row_index, system in enumerate(systems, start=1):
        values = [
            system["label"],
            metric_text(system, "recall"),
            metric_text(system, "correction"),
            metric_text(system, "forgetting"),
            metric_text(system, "quality"),
            str(system["sampleRuns"]),
        ]
        for column, value in enumerate(values):
            cell = table.cell(row_index, column)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = SURFACE if system["label"] == "D_volta_memory" else NAVY
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(11)
            paragraph.font.color.rgb = WHITE if system["label"] == "D_volta_memory" else MUTED
            paragraph.font.bold = system["label"] == "D_volta_memory"
            paragraph.alignment = PP_ALIGN.CENTER


def build_deck(summary: dict, output: Path) -> None:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    blank = presentation.slide_layouts[6]

    # 1 — the claim
    slide = presentation.slides.add_slide(blank)
    add_background(slide)
    add_text(slide, "VOLTA MEMORY", 0.85, 1.6, 11.5, 0.5, size=18, color=GOLD, bold=True)
    add_text(slide, "Persistent AI memory\nthat can prove what changed.", 0.85, 2.15, 9.7, 1.35, size=42, color=WHITE, bold=True)
    add_text(
        slide,
        "A Qwen-powered energy advisor that stores source-linked facts, retains an accountable correction history, and selects only current context for advice.",
        0.85,
        4.05,
        8.4,
        0.8,
        size=20,
        color=MUTED,
    )
    add_text(slide, "Evidence, not assumptions.", 0.85, 5.65, 5.0, 0.35, size=18, color=CYAN, bold=True)
    add_text(slide, "Track 1: MemoryAgent  |  Qwen Cloud Global AI Hackathon", 0.85, 6.85, 8.0, 0.25, size=10, color=MUTED)

    # 2 — lifecycle
    slide = presentation.slides.add_slide(blank)
    add_background(slide)
    add_header(slide, "Memory lifecycle", "Volta keeps the useful fact and preserves the reason it changed.")
    stages = [
        ("1. Extract", "Qwen structures an observation at session end and records its supporting turn."),
        ("2. Verify", "Only a quote matched to a real user message is presented as evidence."),
        ("3. Revise", "A correction supersedes the old fact without deleting its audit history."),
        ("4. Retrieve", "Current, high-value facts are packed to a limited context budget."),
        ("5. Explain", "Every answer distinguishes used evidence from retained or excluded context."),
    ]
    for index, (title, detail) in enumerate(stages):
        x = 0.8 + (index % 3) * 4.05
        y = 1.8 + (index // 3) * 2.2
        add_card(slide, x, y, 3.6, 1.65)
        add_text(slide, title, x + 0.25, y + 0.25, 3.1, 0.28, size=16, color=GOLD, bold=True)
        add_text(slide, detail, x + 0.25, y + 0.68, 3.05, 0.7, size=14, color=MUTED)

    # 3 — demo proof
    slide = presentation.slides.add_slide(blank)
    add_background(slide)
    add_header(slide, "Judge proof sequence", "One correction should be visible from user message to final recommendation.")
    add_card(slide, 0.8, 1.7, 3.6, 4.45)
    add_card(slide, 4.85, 1.7, 3.6, 4.45)
    add_card(slide, 8.9, 1.7, 3.6, 4.45)
    add_text(slide, "1. Confirm", 1.1, 2.05, 3.0, 0.35, size=20, color=GOLD, bold=True)
    add_text(slide, "“My bill is about R3,200, but keeping the lights on is what matters.”", 1.1, 2.7, 2.95, 1.2, size=16, color=WHITE)
    add_text(slide, "2. Correct", 5.15, 2.05, 3.0, 0.35, size=20, color=GOLD, bold=True)
    add_text(slide, "“Actually, the bill is R3,800.”\n\nReceipt: R3,200 → R3,800", 5.15, 2.7, 2.95, 1.2, size=16, color=WHITE)
    add_text(slide, "3. Prove", 9.2, 2.05, 3.0, 0.35, size=20, color=GOLD, bold=True)
    add_text(slide, "The Memory Map shows the active R3,800 fact, its source turn, and the retained R3,200 predecessor.", 9.2, 2.7, 2.95, 1.2, size=16, color=WHITE)
    add_text(slide, "The recommendation must use R3,800 and identify R3,200 as superseded—not merely absent.", 1.1, 5.25, 10.9, 0.45, size=18, color=CYAN, bold=True, align=PP_ALIGN.CENTER)

    # 4 — evaluation
    slide = presentation.slides.add_slide(blank)
    add_background(slide)
    add_header(slide, "Measured trade-offs", "The benchmark reports the full comparison, not only Volta's strongest result.")
    add_metric_table(slide, summary)
    add_text(
        slide,
        f"Run {summary['runId']} · {summary['successfulRuns']}/{summary['expectedRuns']} successful cases · Qwen model {summary['model']}",
        0.8,
        5.8,
        11.75,
        0.25,
        size=12,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )
    add_text(slide, "Read correction, forgetting, recall, latency, and cost together; no single metric is presented as a universal win.", 0.8, 6.25, 11.75, 0.35, size=15, color=CYAN, align=PP_ALIGN.CENTER)

    # 5 — cloud architecture
    slide = presentation.slides.add_slide(blank)
    add_background(slide)
    add_header(slide, "Alibaba Cloud deployment", "Static UI, protected API edge, serverless memory engine.")
    components = [
        ("OSS + CDN", "Static Next.js export\nRuntime API config", 0.8, 2.05, CYAN),
        ("API Gateway", "CORS allowlist\nRate & concurrency limits", 3.45, 2.05, GOLD),
        ("Function Compute 3.0", "FastAPI\nQwen orchestration", 6.1, 2.05, GREEN),
        ("RDS PostgreSQL + pgvector", "Memory lifecycle\nEvidence and relations", 8.75, 2.05, CYAN),
    ]
    for name, detail, x, y, color in components:
        add_card(slide, x, y, 2.25, 1.85)
        add_text(slide, name, x + 0.18, y + 0.28, 1.9, 0.35, size=15, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, detail, x + 0.18, y + 0.78, 1.9, 0.65, size=13, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "Qwen Cloud", 5.25, 4.85, 2.8, 0.32, size=20, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Generation · extraction · relationship classification · embeddings", 3.5, 5.3, 6.3, 0.35, size=16, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, "SLS captures operational telemetry; a least-privilege RAM role accesses only required cloud services.", 1.0, 6.2, 11.3, 0.3, size=14, color=CYAN, align=PP_ALIGN.CENTER)

    # 6 — close
    slide = presentation.slides.add_slide(blank)
    add_background(slide)
    add_text(slide, "THE MEMORY AGENT TEST", 0.85, 1.35, 6.0, 0.25, size=13, color=GOLD, bold=True)
    add_text(slide, "Can the agent show why its recommendation is still true?", 0.85, 1.8, 10.8, 1.0, size=38, color=WHITE, bold=True)
    add_text(slide, "Volta answers with source-linked, current evidence—and keeps the previous fact only for accountability.", 0.85, 3.3, 9.8, 0.6, size=21, color=MUTED)
    app_url = os.environ.get("VOLTA_PUBLIC_APP_ORIGIN", "Configured at release via runtime-config.js")
    add_text(slide, f"Live app: {app_url}", 0.85, 5.15, 9.5, 0.35, size=18, color=CYAN, bold=True)
    add_text(slide, "Repository · benchmark artifact · architecture · public demo video", 0.85, 5.75, 8.5, 0.28, size=15, color=MUTED)
    add_text(slide, "Volta Memory", 0.85, 6.85, 3.2, 0.25, size=10, color=MUTED)

    output.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output)


def main() -> int:
    try:
        build_deck(load_summary(SUMMARY_PATH), OUTPUT_PATH)
    except (FileNotFoundError, json.JSONDecodeError) as exc:
        print(f"Deck generation failed: {exc}", file=sys.stderr)
        return 2
    print(f"Presentation saved to {OUTPUT_PATH}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
